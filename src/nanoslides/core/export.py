"""Slide deck export helpers."""

from __future__ import annotations

from enum import Enum
from pathlib import Path
from typing import Protocol

from PIL import Image
from pptx import Presentation
from pptx.util import Emu

_SUPPORTED_IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".webp"}
_EMU_PER_PIXEL_AT_96_DPI = 9525


class ExportFormat(str, Enum):
    """Supported export targets."""

    PPTX = "pptx"


class SlideDeckExporter(Protocol):
    """Exporter contract for slide deck formats."""

    def export(self, slide_images: list[Path], output_path: Path) -> None:
        """Export slide images into a deck file."""


class PptxSlideDeckExporter:
    """PowerPoint exporter implementation."""

    def export(self, slide_images: list[Path], output_path: Path) -> None:
        image_dimensions = [_load_image_dimensions(path) for path in slide_images]
        uniform_dimensions = _get_uniform_dimensions(image_dimensions)
        presentation = Presentation()
        blank_layout = presentation.slide_layouts[6]
        if uniform_dimensions is not None:
            width_px, height_px = uniform_dimensions
            presentation.slide_width = Emu(width_px * _EMU_PER_PIXEL_AT_96_DPI)
            presentation.slide_height = Emu(height_px * _EMU_PER_PIXEL_AT_96_DPI)

        slide_width = presentation.slide_width
        slide_height = presentation.slide_height

        for image_path in slide_images:
            slide = presentation.slides.add_slide(blank_layout)
            if uniform_dimensions is not None:
                slide.shapes.add_picture(
                    str(image_path),
                    Emu(0),
                    Emu(0),
                    width=slide_width,
                    height=slide_height,
                )
                continue

            picture = slide.shapes.add_picture(str(image_path), Emu(0), Emu(0))
            scale = min(slide_width / picture.width, slide_height / picture.height)
            picture.width = Emu(int(picture.width * scale))
            picture.height = Emu(int(picture.height * scale))
            picture.left = Emu(int((slide_width - picture.width) / 2))
            picture.top = Emu(int((slide_height - picture.height) / 2))

        output_path.parent.mkdir(parents=True, exist_ok=True)
        presentation.save(str(output_path))


_EXPORTERS: dict[ExportFormat, SlideDeckExporter] = {
    ExportFormat.PPTX: PptxSlideDeckExporter(),
}


def _load_image_dimensions(image_path: Path) -> tuple[int, int]:
    try:
        with Image.open(image_path) as image:
            return image.size
    except OSError as exc:
        raise ValueError(f"Unable to read image dimensions from {image_path}.") from exc


def _get_uniform_dimensions(
    dimensions: list[tuple[int, int]],
) -> tuple[int, int] | None:
    first = dimensions[0]
    for size in dimensions[1:]:
        if size != first:
            return None
    return first


def list_slide_images(slides_dir: Path) -> list[Path]:
    """Collect and sort slide image files from a directory."""
    if not slides_dir.exists():
        raise FileNotFoundError(f"Slides directory not found: {slides_dir}")
    if not slides_dir.is_dir():
        raise NotADirectoryError(f"Slides path is not a directory: {slides_dir}")

    slide_images = sorted(
        path
        for path in slides_dir.iterdir()
        if path.is_file() and path.suffix.lower() in _SUPPORTED_IMAGE_SUFFIXES
    )
    if not slide_images:
        raise ValueError(f"No slide images found in {slides_dir}.")
    return slide_images


def export_slides(
    *,
    slides_dir: Path,
    output_path: Path,
    format: ExportFormat = ExportFormat.PPTX,
) -> Path:
    """Export all slide images from a project folder into a deck."""
    slide_images = list_slide_images(slides_dir)
    _EXPORTERS[format].export(slide_images, output_path)
    return output_path
